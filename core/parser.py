import os
import json
import logging
from typing import List, Dict, Any
import io

# Optional imports handled gracefully
try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

try:
    import docx
except ImportError:
    docx = None

try:
    import pandas as pd
except ImportError:
    pd = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    from PIL import Image
    import pytesseract
except ImportError:
    Image = None
    pytesseract = None

logger = logging.getLogger(__name__)

# Dynamically add Tesseract to PATH and set TESSDATA_PREFIX for Windows users
TESSERACT_PATH = r"C:\Program Files\Tesseract-OCR"
if os.path.exists(TESSERACT_PATH):
    os.environ["PATH"] += os.pathsep + TESSERACT_PATH
    os.environ["TESSDATA_PREFIX"] = os.path.join(TESSERACT_PATH, "tessdata")
    if pytesseract:
        pytesseract.pytesseract.tesseract_cmd = os.path.join(TESSERACT_PATH, "tesseract.exe")

def parse_document(file_bytes: bytes, filename: str) -> List[Dict[str, Any]]:
    """
    Master router for parsing documents based on their file extension.
    Supported: pdf, docx, xlsx, csv, txt, jsonl, pptx, png, jpg, jpeg.
    """
    ext = filename.lower().split('.')[-1]
    
    if ext == "pdf":
        return _parse_pdf(file_bytes, filename)
    elif ext == "docx":
        return _parse_docx(file_bytes, filename)
    elif ext in ["xlsx", "csv"]:
        return _parse_excel_csv(file_bytes, filename, ext)
    elif ext == "txt":
        return _parse_txt(file_bytes, filename)
    elif ext == "jsonl":
        return _parse_jsonl(file_bytes, filename)
    elif ext == "pptx":
        return _parse_pptx(file_bytes, filename)
    elif ext in ["png", "jpg", "jpeg"]:
        return _parse_image(file_bytes, filename)
    else:
        raise ValueError(f"Unsupported file format: {ext}")

def _parse_pdf(file_bytes: bytes, filename: str) -> List[Dict[str, Any]]:
    if not fitz:
        raise ImportError("PyMuPDF (fitz) is not installed.")
        
    parsed_pages = []
    try:
        doc = fitz.open(stream=file_bytes, filetype="pdf")
    except Exception as e:
        raise ValueError(f"Could not open PDF file '{filename}': {str(e)}")
        
    for page_idx in range(len(doc)):
        page = doc[page_idx]
        page_num = page_idx + 1
        
        try:
            text = page.get_text().strip()
        except Exception:
            text = ""
            
        if len(text) < 50:
            try:
                textpage = page.get_textpage_ocr(language="eng")
                ocr_text = textpage.extractText().strip()
                if len(ocr_text) > len(text):
                    text = ocr_text
            except Exception as ocr_err:
                logger.error(f"PyMuPDF OCR failed on page {page_num}: {ocr_err}")
                
        if text:
            parsed_pages.append({"text": text, "filename": filename, "page_num": page_num})
            
    doc.close()
    return parsed_pages

def _parse_docx(file_bytes: bytes, filename: str) -> List[Dict[str, Any]]:
    if not docx:
        raise ImportError("python-docx is not installed.")
        
    doc = docx.Document(io.BytesIO(file_bytes))
    full_text = []
    for para in doc.paragraphs:
        if para.text.strip():
            full_text.append(para.text.strip())
            
    text = "\n".join(full_text)
    if text:
        return [{"text": text, "filename": filename, "page_num": 1}]
    return []

def _parse_excel_csv(file_bytes: bytes, filename: str, ext: str) -> List[Dict[str, Any]]:
    if not pd:
        raise ImportError("pandas is not installed.")
        
    parsed_pages = []
    if ext == "csv":
        df = pd.read_csv(io.BytesIO(file_bytes))
        text = df.to_string(index=False)
        if text.strip():
            parsed_pages.append({"text": text, "filename": filename, "page_num": 1})
    else: # xlsx
        # Read all sheets
        xl = pd.ExcelFile(io.BytesIO(file_bytes))
        for sheet_idx, sheet_name in enumerate(xl.sheet_names):
            df = xl.parse(sheet_name)
            text = f"Sheet: {sheet_name}\n" + df.to_string(index=False)
            if text.strip():
                parsed_pages.append({"text": text.strip(), "filename": filename, "page_num": sheet_idx + 1})
                
    return parsed_pages

def _parse_txt(file_bytes: bytes, filename: str) -> List[Dict[str, Any]]:
    text = file_bytes.decode('utf-8', errors='ignore').strip()
    if text:
        return [{"text": text, "filename": filename, "page_num": 1}]
    return []

def _parse_jsonl(file_bytes: bytes, filename: str) -> List[Dict[str, Any]]:
    lines = file_bytes.decode('utf-8', errors='ignore').splitlines()
    text_blocks = []
    for line in lines:
        try:
            data = json.loads(line)
            # Flatten dict to string
            text_blocks.append(json.dumps(data, indent=2))
        except json.JSONDecodeError:
            text_blocks.append(line)
            
    text = "\n\n".join(text_blocks).strip()
    if text:
        return [{"text": text, "filename": filename, "page_num": 1}]
    return []

def _parse_pptx(file_bytes: bytes, filename: str) -> List[Dict[str, Any]]:
    if not Presentation:
        raise ImportError("python-pptx is not installed.")
        
    prs = Presentation(io.BytesIO(file_bytes))
    parsed_pages = []
    for i, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
                
        text = "\n".join(slide_text).strip()
        if text:
            parsed_pages.append({"text": text, "filename": filename, "page_num": i + 1})
            
    return parsed_pages

def _parse_image(file_bytes: bytes, filename: str) -> List[Dict[str, Any]]:
    if not Image or not pytesseract:
        raise ImportError("Pillow or pytesseract is not installed.")
        
    img = Image.open(io.BytesIO(file_bytes))
    text = pytesseract.image_to_string(img).strip()
    
    if text:
        return [{"text": text, "filename": filename, "page_num": 1}]
    return []
